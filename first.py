# -*- coding: utf-8 -*-
"""
Created on Mon Apr 22 08:26:36 2024

@author: u22267
"""


# %m -> Month as a zero-padded decimal number.
# %Y -> Year with century as a decimal number.

import datetime as dt
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
import os

tod = dt.date.today()

today = tod - dt.timedelta(15)

print(today.strftime("%Y-%m-%d %H:%M"))

print("Period :" + today.strftime("%m"))
print("Year :" + today.strftime("%Y"))

ppt_file_name = "HBI-UST_Monthly Operating Report P" + today.strftime("%m") + " v" + "1.0" + ".pptx"
title = "Monthly Operating Report -\n Period " + today.strftime("%m") + " " + today.strftime("%B") + " " + today.strftime("%Y")
parent_dir = "C:\Krishna\Work\MOR"
period = os.path.join(parent_dir, "P" + today.strftime("%m") + "-" + today.strftime("%Y"))

print(ppt_file_name)
print(title)
print(period)

if not os.path.exists(period):
    os.mkdir(period)
    snow = os.path.join(period, "ServiceNow_Reports")
    os.mkdir(snow)
    template_ppt_path = os.path.join(parent_dir, "last_month", "HBI-UST_Monthly Operating Report Template.pptx")
    this_ppt_path = os.path.join(period, ppt_file_name)
    prs = Presentation(template_ppt_path)
    slide1 = prs.slides[0]
    shapes = prs.slides[0].shapes[5]
    slide4 = prs.slides[4]
    # Chart 1 :Incidents Open and Closed in This Month
    chart1 = prs.slides[4].shapes[8].chart
#    for slide in prs.slides:
#        for shape in slide.shapes:
#            print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))
    # ---define new chart data---
    chart_data = CategoryChartData()
    chart_data.categories = ['Open', 'Closed']
    chart_data.add_series('Count', (2000, 1000))
    chart1.replace_data(chart_data)
    # Chart 2 : Ticket Split Based On Priority
    chart2 = prs.slides[4].shapes[9].chart
    chart_data = CategoryChartData()
    chart_data.categories = ['P1', 'P2', 'P3', 'P4','P5']
    chart_data.add_series('Incident', (0, 256, 128, 400, 120))
    chart2.replace_data(chart_data)
    # Chart 3 : Ticket Split Based on Category
    chart3 = prs.slides[4].shapes[10].chart
    chart_data = CategoryChartData()
    chart_data.categories = ['BI', 'Data Management', 'Supply Chain']
    chart_data.add_series('Incident', (656, 66, 561))
    chart3.replace_data(chart_data)
#    for shape in slide4.shapes:
#        print('%s' % shape.shape_type)
    text_frame = shapes.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(25)
    run.font.bold = True
    prs.save(this_ppt_path)
else:
    print("Directory exists !!!!")
        
    
